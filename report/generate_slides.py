# -*- coding: utf-8 -*-
"""
Generates a PowerPoint deck for the CEC418 project
"AI-Based Fraud Detection in School Fee Payments".

Reuses the screenshots in report/images/ (same filenames as the report).
Missing images render as a labelled placeholder. Re-run to rebuild.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
try:
    from PIL import Image
except Exception:
    Image = None

HERE = os.path.dirname(os.path.abspath(__file__))
IMG_DIR = os.path.join(HERE, "images")
OUT = os.path.join(HERE, "Fraud_Detection_CEC418_Presentation.pptx")

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
BLUE = RGBColor(0x2E, 0x6D, 0xA4)
GREY = RGBColor(0x5A, 0x5A, 0x5A)
LGREY = RGBColor(0xEC, 0xEF, 0xF3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x11, 0x11, 0x11)
FONT = "Calibri"
REPO = "github.com/Perry-Bradley/fraud-detection"

IMAGES = {
    "docker": "docker_desktop.png",
    "swagger": "swagger_docs.png",
    "fraud_notify": "fraud_notifications.png",
    "fraud_trend": "fraud_trend.png",
    "circleci": "ci_circleci.png",
    "railway": "railway_production.png",
    "dash_kpi": "dashboard_kpis.png",
    "dash_def": "dashboard_defaulters.png",
}

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _set(run, size, color=BLACK, bold=False, italic=False, font=FONT):
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic; run.font.name = font


def rect(slide, l, t, w, h, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if not line:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    return tf


def slide_number(slide, n):
    tf = textbox(slide, SW - Inches(1.0), SH - Inches(0.5), Inches(0.8), Inches(0.35))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    _set(p.add_run(), 12); p.runs[0].text = str(n); p.runs[0].font.color.rgb = GREY


def footer(slide):
    tf = textbox(slide, Inches(0.5), SH - Inches(0.5), Inches(9), Inches(0.35))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "AI-Based Fraud Detection in School Fee Payments  ·  CEC418"
    _set(r, 9.5, GREY)


def header(slide, title, kicker=None):
    rect(slide, 0, 0, SW, Inches(1.15), NAVY)
    rect(slide, 0, Inches(1.15), SW, Inches(0.06), BLUE)
    tf = textbox(slide, Inches(0.55), Inches(0.18), SW - Inches(1.1), Inches(0.9),
                 anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        p = tf.paragraphs[0]; _set(p.add_run(), 12, RGBColor(0xBF, 0xD3, 0xE6), bold=True)
        p.runs[0].text = kicker.upper()
        p2 = tf.add_paragraph()
    else:
        p2 = tf.paragraphs[0]
    _set(p2.add_run(), 26, WHITE, bold=True); p2.runs[0].text = title


def content_slide(n, title, bullets, kicker=None, two_col=False):
    s = prs.slides.add_slide(BLANK)
    header(s, title, kicker)
    if not two_col:
        tf = textbox(s, Inches(0.7), Inches(1.55), SW - Inches(1.4), SH - Inches(2.3))
        _fill_bullets(tf, bullets)
    else:
        half = (SW - Inches(1.6)) / 2
        tfL = textbox(s, Inches(0.7), Inches(1.55), half, SH - Inches(2.3))
        tfR = textbox(s, Inches(0.9) + half, Inches(1.55), half, SH - Inches(2.3))
        mid = (len(bullets) + 1) // 2
        _fill_bullets(tfL, bullets[:mid]); _fill_bullets(tfR, bullets[mid:])
    footer(s); slide_number(s, n)
    return s


def _fill_bullets(tf, bullets):
    first = True
    for item in bullets:
        if isinstance(item, tuple):
            lead, rest, lvl = item
        else:
            lead, rest, lvl = None, item, 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(8); p.line_spacing = 1.05
        # bullet glyph
        b = p.add_run(); b.text = ("•  " if lvl == 0 else "–  ")
        _set(b, 18 if lvl == 0 else 15, BLUE if lvl == 0 else GREY, bold=True)
        if lead:
            r = p.add_run(); r.text = lead + " "; _set(r, 18 if lvl == 0 else 15, NAVY, bold=True)
        r2 = p.add_run(); r2.text = rest
        _set(r2, 18 if lvl == 0 else 15, BLACK)


def image_fit(slide, path, l, t, w, h):
    if os.path.exists(path) and Image is not None:
        try:
            with Image.open(path) as im:
                iw, ih = im.size
            ar = ih / iw
            box_ar = h / w
            if ar > box_ar:
                ph = h; pw = Emu(int(h / ar))
            else:
                pw = w; ph = Emu(int(w * ar))
            left = l + Emu(int((w - pw) / 2)); top = t + Emu(int((h - ph) / 2))
            slide.shapes.add_picture(path, left, top, width=pw, height=ph)
            return True
        except Exception:
            pass
    elif os.path.exists(path):
        slide.shapes.add_picture(path, l, t, width=w)
        return True
    # placeholder
    box = rect(slide, l, t, w, h, LGREY)
    box.line.color.rgb = GREY; box.line.width = Pt(1)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _set(p.add_run(), 14, NAVY, bold=True); p.runs[0].text = "IMAGE PLACEHOLDER"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    _set(p2.add_run(), 11, GREY, italic=True)
    p2.runs[0].text = f"images/{os.path.basename(path)}"
    return False


def caption(slide, text, l, t, w):
    tf = textbox(slide, l, t, w, Inches(0.5))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _set(p.add_run(), 11, GREY, italic=True); p.runs[0].text = text


def image_slide(n, title, key, capt, kicker=None):
    s = prs.slides.add_slide(BLANK)
    header(s, title, kicker)
    image_fit(s, os.path.join(IMG_DIR, IMAGES[key]),
              Inches(1.3), Inches(1.5), SW - Inches(2.6), Inches(4.7))
    caption(s, capt, Inches(1.3), Inches(6.25), SW - Inches(2.6))
    footer(s); slide_number(s, n)
    return s


def two_image_slide(n, title, key1, key2, kicker=None):
    s = prs.slides.add_slide(BLANK)
    header(s, title, kicker)
    half = (SW - Inches(2.0)) / 2
    image_fit(s, os.path.join(IMG_DIR, IMAGES[key1]), Inches(0.7), Inches(1.6), half, Inches(4.6))
    image_fit(s, os.path.join(IMG_DIR, IMAGES[key2]), Inches(1.3) + half, Inches(1.6), half, Inches(4.6))
    footer(s); slide_number(s, n)
    return s


# =============================================================================
# SLIDE 1 — TITLE
# =============================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.55), SW, Inches(0.05), BLUE)
tf = textbox(s, Inches(0.9), Inches(0.7), SW - Inches(1.8), Inches(0.9))
p = tf.paragraphs[0]; _set(p.add_run(), 15, RGBColor(0xBF, 0xD3, 0xE6), bold=True)
p.runs[0].text = "UNIVERSITY OF BUEA · COLLEGE OF TECHNOLOGY · COMPUTER ENGINEERING"
p2 = tf.add_paragraph(); _set(p2.add_run(), 13, RGBColor(0x9F, 0xB8, 0xD0))
p2.runs[0].text = "CEC418 — Software Construction and Evolution"

tf = textbox(s, Inches(0.9), Inches(2.0), SW - Inches(1.8), Inches(2.4), anchor=MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]; _set(p.add_run(), 44, WHITE, bold=True)
p.runs[0].text = "AI-Based Fraud Detection in School Fee Payments"
p2 = tf.add_paragraph(); p2.space_before = Pt(10)
_set(p2.add_run(), 19, RGBColor(0xCF, 0xDD, 0xEA), italic=True)
p2.runs[0].text = "A machine-learning module built within a School Management System"

tf = textbox(s, Inches(0.9), Inches(4.9), SW - Inches(1.8), Inches(2.0))
for label, val in [("Presented by:", "Sepo Perry-Bradley Dinga  (CT23A145, BTech 400)"),
                   ("Course Instructor:", "Mr. Kometa Denis"),
                   ("Academic Year:", "2025 / 2026"),
                   ("Repository:", REPO)]:
    p = tf.add_paragraph() if tf.paragraphs[0].runs else tf.paragraphs[0]
    r = p.add_run(); r.text = label + "  "; _set(r, 15, RGBColor(0xBF, 0xD3, 0xE6), bold=True)
    r2 = p.add_run(); r2.text = val; _set(r2, 15, WHITE)
    p.space_after = Pt(3)

# =============================================================================
# SLIDE 2 — AGENDA
# =============================================================================
content_slide(2, "Presentation Outline", [
    ("1.", "Introduction & the problem of fee-payment fraud", 0),
    ("2.", "Objectives and how the project maps to CEC418", 0),
    ("3.", "System architecture & technology stack", 0),
    ("4.", "Fraud Detection — the machine-learning core (focus)", 0),
    ("5.", "The DevOps toolchain (focus)", 0),
    ("6.", "Verification, results & evolution", 0),
    ("7.", "Conclusion & future work", 0),
], kicker="Agenda")

# =============================================================================
# SLIDE 3 — INTRODUCTION / PROBLEM
# =============================================================================
content_slide(3, "Introduction & Problem", [
    ("Manual collection:", "schools collect large sums in fees each term, recorded by hand in ledgers or spreadsheets.", 0),
    ("Weak controls:", "balances are hard to reconcile and receipts can be lost or forged.", 0),
    ("Fraud risk:", "the bursary handles cash directly — a natural point for duplicate receipts, altered amounts or unusual patterns.", 0),
    ("The core problem:", "automatically and reliably detect fraudulent or anomalous fee payments the moment they are recorded — without slowing down collection.", 0),
], kicker="Section 1")

# =============================================================================
# SLIDE 4 — OBJECTIVES
# =============================================================================
content_slide(4, "Aim & Objectives", [
    ("", "Record fee payments and auto-generate unique PDF receipts.", 0),
    ("", "Compute real-time student balances and outstanding amounts.", 0),
    ("", "Accept online mobile-money payments (Campay, sandbox/stub mode).", 0),
    ("", "Score every payment with an ML anomaly detector and surface flagged payments for review.", 0),
    ("", "Provide staff dashboards, reporting and an immutable audit trail.", 0),
    ("", "Demonstrate the CEC418 design → construct → test cycle and a full DevOps toolchain.", 0),
], kicker="Section 1")

# =============================================================================
# SLIDE 5 — CEC418 MAPPING
# =============================================================================
content_slide(5, "How It Relates to CEC418", [
    ("Minimising complexity:", "three small single-responsibility services; the ML model isolated in its own service.", 0),
    ("Anticipating change:", "12-factor config; per-environment Helm values; images tagged per commit.", 0),
    ("Constructing for verification:", "unit tests, an immutable audit log, real-time ML scoring, CodeQL/Trivy scanning.", 0),
    ("Reuse & standards:", "OpenAPI schema, a reusable Helm chart, conventional REST.", 0),
    ("Evolution & managing construction:", "Dependabot, GitOps delivery, IaC provisioning, full test re-run per build.", 0),
], kicker="Section 2")

# =============================================================================
# SLIDE 6 — ARCHITECTURE (image)
# =============================================================================
image_slide(6, "System Architecture", "docker",
            "Figure 1 — the full stack running under Docker Compose: backend, ML service, frontend, Postgres, Redis and the observability stack.",
            kicker="Section 3")

# =============================================================================
# SLIDE 7 — ARCHITECTURE DETAIL + STACK
# =============================================================================
content_slide(7, "Services & Technology Stack", [
    ("Backend (Django/DRF, :5000):", "auth, students, fees, payments, receipts, audit, reporting.", 0),
    ("ML service (FastAPI, :8000):", "the fraud-detection model as a stateless scoring endpoint.", 0),
    ("Frontend (React/Vite, :3000):", "a staff console and a student/parent portal behind one login.", 0),
    ("Data tier:", "PostgreSQL 15 for records; Redis 7 for cache and the Celery task queue.", 0),
    ("Self-documenting API:", "drf-spectacular generates the OpenAPI 3.0 / Swagger docs (Figure 2).", 0),
], kicker="Section 3")

# =============================================================================
# SLIDE 8 — FRAUD DETECTION: why + model
# =============================================================================
content_slide(8, "Fraud Detection — The ML Core", [
    ("Real-time scoring:", "every payment is scored on creation via POST /detect-anomaly → { is_anomalous, score, reason }.", 0),
    ("Why a separate service:", "keeps the heavy scientific stack out of the web backend; scales & fails independently.", 0),
    ("Hybrid model:", "adapts to how much history a student has —", 0),
    ("Cold-start heuristic:", "< 10 payments → rule-based (same-day duplicate, z-score > 3; first payment never flagged).", 1),
    ("Isolation Forest:", "≥ 10 payments → scikit-learn IsolationForest, trained leave-one-out, reproducible seed.", 1),
], kicker="Section 4 · Focus")

# =============================================================================
# SLIDE 9 — FRAUD DETECTION: features + resilience
# =============================================================================
content_slide(9, "Features, Scoring & Resilience", [
    ("Five interpretable features:", "amount, amount/mean ratio, z-score, history length, average gap (days).", 0),
    ("Score:", "normalised to 0–1; flagged payments carry a human-readable reason (e.g. “Duplicate amount detected”).", 0),
    ("One path, many sources:", "in-person, online webhook, admin simulate and a Celery task all reuse the same scoring → identical audit logs.", 0),
    ("Fail-open design:", "3-second timeout; if the model is down the payment is treated as not anomalous — a cashier is never blocked.", 0),
], kicker="Section 4 · Focus")

# =============================================================================
# SLIDE 10 — FRAUD DETECTION IN ACTION (2 images)
# =============================================================================
two_image_slide(10, "Fraud Detection in Action", "fraud_notify", "fraud_trend",
                kicker="Section 4 · Focus")
# add a caption strip
_s = prs.slides[-1]
caption(_s, "Figures 3–4 — live anomaly notifications (left) and the detection trend with the most common anomaly reasons (right).",
        Inches(0.7), Inches(6.35), SW - Inches(1.4))

# =============================================================================
# SLIDE 11 — DEVOPS OVERVIEW
# =============================================================================
content_slide(11, "The DevOps Toolchain", [
    ("Source / build:", "Git + GitHub · Docker (multi-stage, non-root) · Docker Compose for the local stack.", 0),
    ("CI / CD:", "CircleCI (active) runs all tests; GitHub Actions mirror builds, scans & publishes images.", 0),
    ("Security:", "Trivy image CVE scanning · CodeQL SAST · Dependabot dependency PRs.", 0),
    ("Infra & orchestration:", "Terraform (IaC) · Ansible (config) · Kubernetes · Helm · Argo CD (GitOps).", 0),
    ("Secrets & observability:", "HashiCorp Vault · Prometheus + Grafana + Loki/Promtail.", 0),
], kicker="Section 5 · Focus")

# =============================================================================
# SLIDE 12 — CI/CD pipeline (image)
# =============================================================================
image_slide(12, "CI/CD — Automated Pipeline", "circleci",
            "Figure 5 — the CircleCI pipeline (ml-tests, frontend-build, backend-tests, deploy). CircleCI was chosen to demonstrate a tool beyond GitHub's own ecosystem.",
            kicker="Section 5 · Focus")

# =============================================================================
# SLIDE 13 — PRODUCTION DEPLOYMENT (Railway)
# =============================================================================
image_slide(13, "Live Production Deployment (Railway)", "railway",
            "Figure 6 — the production environment on Railway: backend, ML fraud-detection and frontend services plus managed Postgres & Redis, all online, auto-deployed from GitHub on every push to main.",
            kicker="Section 5 · Focus")

# =============================================================================
# SLIDE 14 — DEVOPS FLOW
# =============================================================================
content_slide(14, "From Commit to Cloud — Automatically", [
    ("1.", "Developer pushes to GitHub.", 0),
    ("2.", "CircleCI runs the full test suite; images are built and scanned with Trivy.", 0),
    ("3.", "Passing images are published to the registry, tagged by commit SHA.", 0),
    ("4.", "Argo CD syncs the updated Helm chart onto the Terraform-provisioned, Ansible-configured Kubernetes cluster.", 0),
    ("5.", "Services emit metrics & logs to Prometheus / Loki / Grafana; Vault supplies secrets.", 0),
    ("Result:", "every push to main reaches a running, observed deployment with no manual step.", 0),
], kicker="Section 5 · Focus")

# =============================================================================
# SLIDE 14 — RESULTS (2 images)
# =============================================================================
two_image_slide(15, "Verification & Results", "dash_kpi", "dash_def", kicker="Section 6")
_s = prs.slides[-1]
caption(_s, "Figures 7–8 — management dashboards: collection KPIs and the AI anomaly count (left); top defaulters and live activity feed (right).",
        Inches(0.7), Inches(6.35), SW - Inches(1.4))

# =============================================================================
# SLIDE 15 — VERIFICATION + EVOLUTION
# =============================================================================
content_slide(16, "Verification & Evolution", [
    ("Backend tests:", "verify receipt numbers are auto-generated and unique.", 0),
    ("ML tests:", "first payment never flagged; same-day duplicate flagged; consistent history not flagged.", 0),
    ("Continuous checks:", "CodeQL, Trivy and Dependabot guard code quality and the supply chain.", 0),
    ("Built to evolve:", "grew from a fee tracker into a full school platform without destabilising the core.", 0),
    ("Traceability:", "commit-tagged images + Argo CD revision history → any behaviour can be traced and rolled back.", 0),
], kicker="Section 6")

# =============================================================================
# SLIDE 16 — CONCLUSION
# =============================================================================
content_slide(17, "Conclusion & Future Work", [
    ("Delivered:", "a real-time ML fraud-detection capability for school-fee payments, inside a complete DevOps toolchain.", 0),
    ("Embodies CEC418:", "minimising complexity, anticipating change, constructing for verification, reuse & standards.", 0),
    ("Future work:", "", 0),
    ("", "a settled production payment provider;", 1),
    ("", "retrain the model on real data with periodic evaluation;", 1),
    ("", "higher test coverage and Prometheus alerting.", 1),
], kicker="Section 7")

# =============================================================================
# SLIDE 17 — THANK YOU
# =============================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.2), SW, Inches(0.05), BLUE)
tf = textbox(s, Inches(0.9), Inches(2.6), SW - Inches(1.8), Inches(1.6), anchor=MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
_set(p.add_run(), 48, WHITE, bold=True); p.runs[0].text = "Thank You"
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(8)
_set(p2.add_run(), 20, RGBColor(0xCF, 0xDD, 0xEA), italic=True); p2.runs[0].text = "Questions & Discussion"
tf = textbox(s, Inches(0.9), Inches(4.6), SW - Inches(1.8), Inches(1.0))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Repository:  " + REPO; _set(r, 16, WHITE)
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r = p2.add_run(); r.text = "Sepo Perry-Bradley Dinga · CT23A145 · CEC418"; _set(r, 14, RGBColor(0x9F, 0xB8, 0xD0))

prs.save(OUT)
present = sum(1 for f in IMAGES.values() if os.path.exists(os.path.join(IMG_DIR, f)))
print("Saved:", OUT)
print("Slides:", len(prs.slides.__iter__.__self__._sldIdLst), "| images embedded:", present, "/", len(IMAGES))
